import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

def extract_text_from_pdf(pdf_path):
    text = ""
    doc = fitz.open(pdf_path)
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_all_course_materials_text(course_materials_dir):
    all_texts = {}
    for root, _, files in os.walk(course_materials_dir):
        for file in files:
            path = os.path.join(root, file)
            if file.lower().endswith(".pdf"):
                all_texts[file] = extract_text_from_pdf(path)
            elif file.lower().endswith(".docx"):
                all_texts[file] = extract_text_from_docx(path)
            elif file.lower().endswith(".pptx"):
                all_texts[file] = extract_text_from_pptx(path)
    return all_texts

if __name__ == "__main__":
    course_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "course_materials"))
    texts = extract_all_course_materials_text(course_dir)
    for filename, text in texts.items():
        print(f"Extracted text from {filename} (length={len(text)})")

    # Save extracted text as JSON files in extracted_text folder
    import json
    extracted_text_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "extracted_text"))
    os.makedirs(extracted_text_dir, exist_ok=True)

    for filename, text in texts.items():
        base_name = os.path.splitext(filename)[0]
        json_path = os.path.join(extracted_text_dir, base_name + ".json")
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump({"filename": filename, "text": text}, f, ensure_ascii=False, indent=2)
